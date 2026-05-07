from __future__ import annotations

from pathlib import Path
from textwrap import dedent

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from .config import DOCS_DIR, FIGURES_DIR, METRICS_DIR, PAPER_DIR, POSTER_DIR


def _read_csv(name: str) -> pd.DataFrame | None:
    path = METRICS_DIR / name
    if not path.exists():
        return None
    return pd.read_csv(path)


def _fmt_pct(value: float | None) -> str:
    if value is None or pd.isna(value):
        return "pending"
    return f"{value * 100:.1f}%"


def _fmt_num(value: float | None, digits: int = 2) -> str:
    if value is None or pd.isna(value):
        return "pending"
    return f"{value:.{digits}f}"


def summarize_results() -> dict[str, str]:
    summary = _read_csv("summary.csv")
    benchmark = _read_csv("benchmark_seq_len.csv")
    noise = _read_csv("noise_sweep.csv")
    gradient = _read_csv("gradient_variance.csv")
    alignment = _read_csv("attention_alignment.csv")

    values: dict[str, str] = {}
    if summary is not None:
        by_model = {row.model_type: row for row in summary.itertuples()}
        values["classical_acc"] = _fmt_pct(getattr(by_model.get("classical"), "test_accuracy", None))
        values["hybrid_acc"] = _fmt_pct(getattr(by_model.get("hybrid_quantum"), "test_accuracy", None))
        values["ablation_acc"] = _fmt_pct(getattr(by_model.get("classical_ablation"), "test_accuracy", None))
        best = summary.sort_values("test_accuracy", ascending=False).iloc[0]
        values["best_model"] = str(best["model_type"])
        values["best_acc"] = _fmt_pct(float(best["test_accuracy"]))
    else:
        values.update({"classical_acc": "pending", "hybrid_acc": "pending", "ablation_acc": "pending", "best_model": "pending", "best_acc": "pending"})

    if benchmark is not None:
        sequence_values = sorted(int(value) for value in benchmark["sequence_length"].unique())
        seq = 32 if 32 in sequence_values else sequence_values[len(sequence_values) // 2]
        subset = benchmark[benchmark["sequence_length"] == seq]
        by_model = {row.model_type: row for row in subset.itertuples()}
        classical = getattr(by_model.get("classical"), "mean_forward_ms", None)
        hybrid = getattr(by_model.get("hybrid_quantum"), "mean_forward_ms", None)
        ratio = float(hybrid) / float(classical) if classical and hybrid else None
        values["runtime_seq"] = str(seq)
        values["hybrid_runtime_ratio"] = _fmt_num(ratio, 2)
    else:
        values["runtime_seq"] = "pending"
        values["hybrid_runtime_ratio"] = "pending"

    if noise is not None and "hybrid_quantum" in noise["model_type"].values:
        h = noise[noise["model_type"] == "hybrid_quantum"].sort_values("noise_level")
        start_acc = float(h.iloc[0]["test_accuracy"])
        end_acc = float(h.iloc[-1]["test_accuracy"])
        delta = end_acc - start_acc
        values["hybrid_noise_change"] = f"{delta * 100:+.1f} percentage points ({_fmt_pct(start_acc)} to {_fmt_pct(end_acc)})"
    else:
        values["hybrid_noise_change"] = "pending"

    if gradient is not None and len(gradient) > 1:
        first = float(gradient.sort_values("depth").iloc[0]["grad_variance"])
        last = float(gradient.sort_values("depth").iloc[-1]["grad_variance"])
        trend = "decreased" if last < first else "increased"
        values["gradient_trend"] = trend
        values["gradient_start"] = _fmt_num(first, 6)
        values["gradient_end"] = _fmt_num(last, 6)
    else:
        values["gradient_trend"] = "pending"
        values["gradient_start"] = "pending"
        values["gradient_end"] = "pending"

    if alignment is not None:
        row = alignment[(alignment["left_model"] == "classical") & (alignment["right_model"] == "hybrid_quantum")]
        values["cka_classical_hybrid"] = _fmt_num(float(row.iloc[0]["linear_cka"]), 3) if not row.empty else "pending"
    else:
        values["cka_classical_hybrid"] = "pending"
    return values


def write_abstract() -> Path:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    r = summarize_results()
    final_abstract = dedent(
        f"""
        # Abstract

        Large language models rely on self-attention, but the query-key similarity step scales quadratically with sequence length. This project tested a compact hybrid quantum-classical attention block that replaces only the classical query-key dot product with a simulated four-qubit parameterized quantum circuit while preserving classical value aggregation, residual structure, and classification layers. Three matched AG News classifiers were implemented: a standard scaled dot-product attention model, the proposed hybrid quantum-kernel model, and a low-dimensional classical ablation using the same query/key bottleneck as the quantum version. The study evaluated classification quality, runtime scaling, attention-map alignment, gradient variance across circuit depth, and robustness under simulated quantum noise. In the compact run, the classical, hybrid, and ablation models reached {r["classical_acc"]}, {r["hybrid_acc"]}, and {r["ablation_acc"]} test accuracy, respectively. At sequence length {r["runtime_seq"]}, the hybrid forward pass was {r["hybrid_runtime_ratio"]}x the classical runtime, showing that state-vector simulation does not provide practical speedup on classical hardware. Attention alignment between classical and hybrid maps had CKA {r["cka_classical_hybrid"]}, and the gradient-variance diagnostic {r["gradient_trend"]} from {r["gradient_start"]} to {r["gradient_end"]} across tested depths. Overall, the experiment supports the feasibility of isolating quantum similarity inside attention, but frames the result as a diagnostic study rather than evidence of immediate efficiency advantage.

        ## Poster Abstract

        This project implements a hybrid quantum-classical attention block for AG News classification. A four-qubit simulated circuit replaces only the query-key similarity calculation, while the value path and classifier remain classical. Matched classical and low-dimensional ablation baselines show whether the quantum kernel changes attention behavior, trainability, robustness, or runtime. The compact results indicate feasible integration but no classical-simulation speedup, making the project most useful as an honest diagnostic of where quantum attention helps and where it remains constrained.
        """
    ).strip()
    out = DOCS_DIR / "abstract.md"
    out.write_text(final_abstract + "\n", encoding="utf-8")
    return out


def write_paper_markdown() -> Path:
    PAPER_DIR.mkdir(parents=True, exist_ok=True)
    r = summarize_results()
    paper = dedent(
        f"""
        # A Hybrid Quantum-Classical Attention Mechanism for Efficient Large Language Models

        **Authors:** Ansh Malhotra and Smaran Kudapa

        ## Abstract

        Large language models rely on self-attention, but query-key similarity scales quadratically with sequence length. This project implements a compact hybrid quantum-classical attention block that replaces only the query-key dot product with a simulated four-qubit quantum kernel. On AG News, the classical, hybrid, and matched low-dimensional ablation models reached {r["classical_acc"]}, {r["hybrid_acc"]}, and {r["ablation_acc"]} test accuracy. The hybrid simulation was {r["hybrid_runtime_ratio"]}x the classical runtime at sequence length {r["runtime_seq"]}, so the result does not support a practical classical-simulation speedup. Instead, the project demonstrates a reproducible framework for isolating quantum similarity inside attention and measuring expressivity, trainability, efficiency, and noise robustness.

        ## Introduction

        Transformer models use self-attention to compare every token with every other token. This design is powerful, but the query-key similarity matrix requires O(n^2) pairwise interactions as sequence length grows. Quantum machine learning suggests a possible alternative: encode low-dimensional query and key features into Hilbert space and compute similarity through state overlap. Prior quantum attention work, including QSANN and QMSAN, shows that quantum self-attention can be applied to text classification, but those models do not fully isolate the query-key similarity step from other architectural changes.

        ## Research Question

        Under controlled conditions, does replacing the query-key dot product with a simulated quantum kernel improve attention expressivity, preserve trainability, and reduce empirical cost?

        ## Methods

        The experiment uses AG News classification with a simple local tokenizer, fixed vocabulary, and deterministic train/validation/test subsets. Three matched PyTorch models are compared. The classical model uses scaled dot-product attention. The hybrid model projects query and key vectors into four dimensions, encodes them into a four-qubit state-vector circuit using repeated RY rotations and nearest-neighbor CZ gates, then computes similarity as |<psi(q)|psi(k)>|^2. The ablation model uses the same four-dimensional query/key bottleneck but computes dot-product similarity classically.

        Evaluation includes classification accuracy and macro-F1, forward-pass runtime by sequence length, memory deltas, centered kernel alignment (CKA) between attention maps, gradient variance across circuit depth, and simulated quantum noise. The noise test applies angle perturbation and depolarizing-style overlap mixing to the hybrid circuit.

        ## Results

        The best compact-run classifier was **{r["best_model"]}** with {r["best_acc"]} test accuracy. The classical baseline reached {r["classical_acc"]}, the hybrid model reached {r["hybrid_acc"]}, and the ablation reached {r["ablation_acc"]}. At sequence length {r["runtime_seq"]}, the hybrid forward pass was {r["hybrid_runtime_ratio"]}x the classical runtime, which is expected because a state-vector circuit is being simulated on classical hardware. The classical-hybrid attention CKA was {r["cka_classical_hybrid"]}. The gradient variance diagnostic {r["gradient_trend"]} from {r["gradient_start"]} to {r["gradient_end"]} over the tested circuit depths. At the highest simulated noise level, hybrid accuracy changed by {r["hybrid_noise_change"]}.

        ## Discussion

        The main finding is not that the simulated quantum layer is faster; it is not. The useful result is that the project isolates the quantum kernel inside an otherwise classical attention block and measures the consequences directly. If the hybrid model performs similarly to the bottlenecked ablation, then most of the behavior comes from dimensionality reduction rather than quantum geometry. If the hybrid model differs in attention alignment, noise response, or gradient behavior, those diagnostics identify where quantum kernels may be worth studying further.

        ## Limitations

        This project uses a small classifier and a state-vector simulator, not quantum hardware. The reported runtime is therefore a measure of classical simulation overhead rather than a claim about future hardware advantage. The dataset is AG News only, and the model is intentionally small enough for reproducible senior-project experiments.

        ## Conclusion

        A hybrid quantum-classical attention block can be implemented end-to-end and evaluated with real metrics. In this compact study, it did not demonstrate practical efficiency gains under classical simulation, but it produced a reproducible framework for testing attention expressivity, trainability, and robustness while keeping the quantum intervention isolated to the query-key similarity operation.

        ## References

        - Vaswani et al. (2017). Attention Is All You Need.
        - Li, Zhao, and Wang (2023). Quantum Self-Attention Neural Networks for Text Classification. arXiv:2205.05625.
        - Chen et al. (2025). Quantum Mixed-State Self-Attention Network. Neural Networks, 185, 107123.
        - Kornblith et al. (2019). Similarity of Neural Network Representations Revisited. ICML.
        - McClean et al. (2018). Barren Plateaus in Quantum Neural Network Training Landscapes. Nature Communications.
        - Shen et al. (2025). Squat: Quant Small Language Models on the Edge. arXiv:2402.10787.
        """
    ).strip()
    out = PAPER_DIR / "final_paper.md"
    out.write_text(paper + "\n", encoding="utf-8")
    return out


def write_paper_pdf() -> Path:
    write_paper_markdown()
    r = summarize_results()
    out = PAPER_DIR / "final_paper.pdf"
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TitleCenter", parent=styles["Title"], alignment=TA_CENTER, fontSize=18, leading=22))
    styles.add(ParagraphStyle(name="Small", parent=styles["BodyText"], fontSize=8.5, leading=11))
    story = [
        Paragraph("A Hybrid Quantum-Classical Attention Mechanism for Efficient Large Language Models", styles["TitleCenter"]),
        Paragraph("Ansh Malhotra and Smaran Kudapa", styles["BodyText"]),
        Spacer(1, 0.16 * inch),
        Paragraph("Abstract", styles["Heading1"]),
        Paragraph(
            f"This compact study replaces only the query-key dot product in attention with a simulated four-qubit quantum kernel. On AG News, the classical, hybrid, and ablation models reached {r['classical_acc']}, {r['hybrid_acc']}, and {r['ablation_acc']} test accuracy. The hybrid model was {r['hybrid_runtime_ratio']}x the classical runtime at sequence length {r['runtime_seq']}, showing no classical-simulation speedup but providing a reproducible diagnostic framework.",
            styles["BodyText"],
        ),
        Spacer(1, 0.12 * inch),
    ]
    image_paths = [
        FIGURES_DIR / "architecture_diagram.png",
        FIGURES_DIR / "accuracy_f1_summary.png",
        FIGURES_DIR / "benchmark_runtime.png",
        FIGURES_DIR / "noise_robustness.png",
        FIGURES_DIR / "gradient_variance.png",
        FIGURES_DIR / "attention_alignment.png",
    ]
    for path in image_paths:
        if path.exists():
            story.append(Image(str(path), width=6.2 * inch, height=3.5 * inch))
            story.append(Spacer(1, 0.12 * inch))
    story.extend(
        [
            Paragraph("Interpretation", styles["Heading1"]),
            Paragraph(
                "The result should be read as a feasibility and diagnostic study. The hybrid circuit can be integrated into attention and compared against matched baselines, but state-vector simulation adds overhead. The strongest contribution is the controlled experiment design: accuracy, runtime, CKA, gradient variance, and noise robustness are all measured from saved metric files.",
                styles["BodyText"],
            ),
            Paragraph("References", styles["Heading1"]),
            Paragraph(
                "Vaswani et al. (2017); Li et al. (2023), arXiv:2205.05625; Chen et al. (2025), Neural Networks 185; Kornblith et al. (2019), ICML; McClean et al. (2018), Nature Communications; Shen et al. (2025), arXiv:2402.10787.",
                styles["Small"],
            ),
        ]
    )
    doc = SimpleDocTemplate(str(out), pagesize=letter, rightMargin=0.65 * inch, leftMargin=0.65 * inch, topMargin=0.65 * inch, bottomMargin=0.65 * inch)
    doc.build(story)
    return out


def _add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=(35, 51, 58), fill=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.color.rgb = RGBColor(200, 210, 214)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return shape


def write_poster_pptx() -> Path:
    POSTER_DIR.mkdir(parents=True, exist_ok=True)
    r = summarize_results()
    prs = Presentation()
    prs.slide_width = Inches(48)
    prs.slide_height = Inches(36)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 249, 246)
    _add_textbox(slide, 1.1, 0.6, 45.8, 1.5, "A Hybrid Quantum-Classical Attention Mechanism for Efficient LLMs", size=32, bold=True)
    _add_textbox(slide, 1.15, 2.0, 32, 0.55, "Ansh Malhotra and Smaran Kudapa | QLab Senior Research Project", size=16, color=(78, 88, 94))
    _add_textbox(slide, 1.1, 3.2, 14.5, 6.0, "Research Question\n\nCan a simulated quantum kernel replace only the query-key dot product while preserving useful attention behavior?\n\nModels\n\nClassical attention\nMatched low-dimensional ablation\nHybrid four-qubit PQC kernel", size=18, fill=(232, 240, 242))
    _add_textbox(slide, 16.4, 3.2, 14.8, 6.0, f"Main Result\n\nClassical accuracy: {r['classical_acc']}\nHybrid accuracy: {r['hybrid_acc']}\nAblation accuracy: {r['ablation_acc']}\n\nAt sequence length {r['runtime_seq']}, the hybrid simulation took {r['hybrid_runtime_ratio']}x the classical runtime.", size=18, fill=(242, 235, 224))
    _add_textbox(slide, 31.9, 3.2, 15.0, 6.0, f"Conclusion\n\nThe hybrid block runs end-to-end, but classical state-vector simulation does not deliver speedup. The value is diagnostic: attention alignment, gradient variance, and noise response can now be measured cleanly.", size=18, fill=(236, 240, 232))

    image_slots = [
        ("architecture_diagram.png", 1.1, 10.2, 22.2, 11.2),
        ("accuracy_f1_summary.png", 24.2, 10.2, 10.9, 8.1),
        ("benchmark_runtime.png", 36.0, 10.2, 10.9, 8.1),
        ("noise_robustness.png", 1.1, 22.3, 14.5, 9.4),
        ("gradient_variance.png", 16.9, 22.3, 14.5, 9.4),
        ("attention_alignment.png", 32.5, 22.3, 14.5, 9.4),
    ]
    for name, x, y, w, h in image_slots:
        path = FIGURES_DIR / name
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    _add_textbox(slide, 1.1, 33.2, 45.8, 0.8, "Real metrics are generated by scripts/run_all.py and saved as CSV files in results/metrics/. Negative or mixed findings are reported directly.", size=15, color=(78, 88, 94))
    out = POSTER_DIR / "QLab_Hybrid_Quantum_Attention_Poster.pptx"
    prs.save(out)
    return out


def write_poster_pdf() -> Path:
    r = summarize_results()
    out = POSTER_DIR / "QLab_Hybrid_Quantum_Attention_Poster.pdf"
    width, height = 48 * inch, 36 * inch
    c = canvas.Canvas(str(out), pagesize=(width, height))
    c.setFillColorRGB(0.972, 0.976, 0.965)
    c.rect(0, 0, width, height, fill=1, stroke=0)

    def xywh(x: float, y_from_top: float, w: float, h: float) -> tuple[float, float, float, float]:
        return x * inch, height - (y_from_top + h) * inch, w * inch, h * inch

    def box(x: float, y: float, w: float, h: float, fill: tuple[float, float, float]) -> None:
        px, py, pw, ph = xywh(x, y, w, h)
        c.setFillColorRGB(*fill)
        c.setStrokeColorRGB(0.82, 0.85, 0.85)
        c.roundRect(px, py, pw, ph, 8, fill=1, stroke=1)

    def wrapped(text: str, x: float, y: float, w: float, font_size: int = 15, bold: bool = False, leading: int | None = None) -> None:
        leading = leading or int(font_size * 1.22)
        px = x * inch
        top = height - y * inch
        max_chars = max(18, int(w * 8.3 * (15 / font_size)))
        c.setFillColorRGB(0.12, 0.16, 0.18)
        c.setFont("Helvetica-Bold" if bold else "Helvetica", font_size)
        text_obj = c.beginText(px, top)
        text_obj.setLeading(leading)
        for paragraph in text.split("\n"):
            words = paragraph.split()
            line = ""
            for word in words:
                candidate = f"{line} {word}".strip()
                if len(candidate) > max_chars and line:
                    text_obj.textLine(line)
                    line = word
                else:
                    line = candidate
            text_obj.textLine(line)
            if not words:
                text_obj.textLine("")
        c.drawText(text_obj)

    c.setFillColorRGB(0.08, 0.10, 0.12)
    c.setFont("Helvetica-Bold", 30)
    c.drawCentredString(width / 2, height - 0.9 * inch, "A Hybrid Quantum-Classical Attention Mechanism for Efficient LLMs")
    c.setFont("Helvetica", 14)
    c.drawString(1.1 * inch, height - 1.75 * inch, "Ansh Malhotra and Smaran Kudapa | QLab Senior Research Project")

    box(1.1, 2.7, 14.6, 4.4, (0.91, 0.94, 0.95))
    box(16.7, 2.7, 14.6, 4.4, (0.95, 0.92, 0.86))
    box(32.3, 2.7, 14.6, 4.4, (0.92, 0.94, 0.91))
    wrapped("Research Question\nCan a simulated quantum kernel replace only query-key similarity while preserving attention behavior?", 1.45, 3.25, 13.8, 15, True)
    wrapped(f"Main Result\nClassical: {r['classical_acc']}\nHybrid: {r['hybrid_acc']}\nAblation: {r['ablation_acc']}\nHybrid runtime ratio: {r['hybrid_runtime_ratio']}x", 17.05, 3.25, 13.8, 15, True)
    wrapped("Conclusion\nThe hybrid layer is feasible, but state-vector simulation is slower. The framework is most useful for diagnostic comparison.", 32.65, 3.25, 13.8, 15, True)

    image_slots = [
        ("architecture_diagram.png", 1.1, 8.0, 21.9, 11.8),
        ("accuracy_f1_summary.png", 23.6, 8.0, 11.3, 11.8),
        ("benchmark_runtime.png", 35.6, 8.0, 11.3, 11.8),
        ("noise_robustness.png", 1.1, 20.8, 14.8, 11.7),
        ("gradient_variance.png", 16.6, 20.8, 14.8, 11.7),
        ("attention_alignment.png", 32.1, 20.8, 14.8, 11.7),
    ]
    for name, x, y, w, h in image_slots:
        path = FIGURES_DIR / name
        if path.exists():
            px, py, pw, ph = xywh(x, y, w, h)
            c.setStrokeColorRGB(0.86, 0.88, 0.88)
            c.rect(px, py, pw, ph, fill=0, stroke=1)
            c.drawImage(str(path), px + 0.15 * inch, py + 0.15 * inch, width=pw - 0.3 * inch, height=ph - 0.3 * inch, preserveAspectRatio=True, anchor="c")
    c.setFont("Helvetica", 12)
    c.setFillColorRGB(0.30, 0.35, 0.37)
    c.drawString(1.1 * inch, 1.0 * inch, "Metrics are generated by scripts/run_all.py and saved as CSV files in results/metrics/. Negative or mixed findings are reported directly.")
    c.showPage()
    c.save()
    return out


def write_all_artifacts() -> list[Path]:
    outputs = [write_abstract(), write_paper_markdown(), write_paper_pdf(), write_poster_pptx(), write_poster_pdf()]
    return outputs
